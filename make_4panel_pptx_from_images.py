# ---
# jupyter:
#   jupytext:
#     formats: ipynb,py:light
#     text_representation:
#       extension: .py
#       format_name: light
#       format_version: '1.5'
#       jupytext_version: 1.4.1
#   kernelspec:
#     display_name: Python 3
#     language: python
#     name: python3
# ---

# +
from pptx import Presentation
from pptx.util import Inches

from datetime import datetime
from skimage import color
import glob
import cv2
import numpy as np
from io import BytesIO
from tqdm.autonotebook import tqdm
import PIL.Image as Image

# +
# -- Set meta data which will appear on first slide
title = "Epi/stroma segmentation"
date = datetime.today()
author = "Andrew Janowczyk"
comments = "data and code taken from blog andrewjanowczyk.com "
pptxfname = "epistroma_results.pptx"

#we only want to generate output for images which have masks, so we find those files
mask_files=glob.glob('./masks/*.png')

# +
#create presentation 
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(10)

blank_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(blank_slide_layout)

#make first slide with our metadata
slide.placeholders[0].text = title

tf = slide.placeholders[1].text_frame
tf.text = f'Date: {date}\n'
tf.text += f"Author: {author}\n"
tf.text += f"Comments: {comments}\n"

# -

#wrapper function to add an image as a byte stream to a slide
#note that this is in place of having to save output directly to disk, and can be used in dynamic settings as well
def addimagetoslide(slide,img,left,top, height, width, resize = .5):
    res = cv2.resize(img , None, fx=resize,fy=resize ,interpolation=cv2.INTER_CUBIC) #since the images are going to be small, we can resize them to prevent the final pptx file from being large for no reason
    image_stream = BytesIO()
    Image.fromarray(res).save(image_stream,format="PNG")

    pic = slide.shapes.add_picture(image_stream, left, top ,height,width)
    image_stream.close()


#helper function to blend two images
def blend2Images(img, mask):
    if (img.ndim == 3):
        img = color.rgb2gray(img)
    if (mask.ndim == 3):
        mask = color.rgb2gray(mask)
    img = img[:, :, None] * 1.0  # can't use boolean
    mask = mask[:, :, None] * 1.0
    out = np.concatenate((mask, img, mask), 2) * 255
    return out.astype('uint8')


# +
for mask_fname in tqdm(mask_files):
    
    #add a new slide for this set of images 
    blank_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_slide_layout)

    
    #compute the associated filenames that we'll need
    orig_fname=mask_fname.replace("./masks","./imgs").replace("_mask.png",".tif")
    output_fname=mask_fname.replace("./masks","./output").replace("_mask.png","_class.png")
    
    
        
    
    #------- orig  - load and add to slide
    img = cv2.cvtColor(cv2.imread(orig_fname),cv2.COLOR_BGR2RGB)
    addimagetoslide(slide, img, Inches(0),Inches(0),Inches(5),Inches(5))
    
    #------ mask - load and add to slide
    mask = cv2.cvtColor(cv2.imread(mask_fname),cv2.COLOR_BGR2RGB)
    addimagetoslide(slide, mask, Inches(5),Inches(0),Inches(5),Inches(5))
    
    #------ output - load and add to slide
    output = cv2.cvtColor(cv2.imread(output_fname),cv2.COLOR_BGR2RGB)
    addimagetoslide(slide, output, Inches(5),Inches(5),Inches(5),Inches(5))
    
    #------ Fuse - load and add to slide
    addimagetoslide(slide,blend2Images(output,mask), Inches(0),Inches(5),Inches(5),Inches(5))
    
    
    #------ Lastly we can also add some metrics/results/values if we would like
    # here we do simple FP/TP/TN/FN
    txBox = slide.shapes.add_textbox(Inches(10), Inches(0),Inches(4),Inches(4) )
    tf = txBox.text_frame
    tf.text = f"{orig_fname}\n"
    tf.text += f"Overall Pixel Agreement: {(output==mask).mean():.4f}\n"
    tf.text += f"True Positive Rate: {(mask[output>0]>0).sum()/(output>0).sum():.4f}\n"
    tf.text += f"False Positive Rate: {(mask[output==0]>0).sum()/(output==0).sum():.4f}\n"
    tf.text += f"True Negative Rate: {(mask[output==0]==0).sum()/(output==0).sum():.4f}\n"
    tf.text += f"False Negative Rate: {(mask[output>0]==0).sum()/(output>0).sum():.4f}\n"
    
    
#At this point the pptx has not been saved, so we do that here and we're all done!
prs.save(pptxfname)
