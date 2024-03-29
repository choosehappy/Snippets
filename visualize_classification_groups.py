# -*- coding: utf-8 -*-
# blogpost - http://www.andrewjanowczyk.com/?p=1661
# Written by Jackson Jacobs - jjj72@case.edu, 2022

import argparse
from PIL import Image
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from tqdm.autonotebook import tqdm
import matplotlib.pyplot as plt
import tables

def addimagetoslide(slide,image_stream,left,top, height, width, resize = 1.0, comment = ''):
	slide.shapes.add_picture(image_stream, left, top ,height,width)
	txBox = slide.shapes.add_textbox(left, Inches(1), width, height)
	tf = txBox.text_frame
	tf.text = comment

if __name__ == "__main__":
	parser = argparse.ArgumentParser()
	parser.add_argument('pytable_path', type=str, help='The path to an h5 file that contains "patch", "ground_truth_label", and "fname" datasets.')	
	parser.add_argument('ppt_save', type=str, help='The full path to save the generated powerpoint.')

	parser.add_argument('--str_criteria', '-s', default=["TRUE NEGATIVES", "FALSE NEGATIVES", "FALSE POSITIVES", "TRUE POSITIVES"], type=str, nargs="+")
	parser.add_argument('--criteria', default=['00', '10', '01', '11'], type=str, nargs="+", help="Each item is a pair of ground-truth,prediction labels.")
	parser.add_argument('--num_rows', '-r',default=5, type=int, help="The number of rows of images per slide.")
	parser.add_argument('--num_cols', '-c',default=8, type=int, help="The number of columns of images per slide.")
	args = parser.parse_args()
	

	##### INITIALIZE CRITERIA DICT #####
	criteria_dict = {args.str_criteria[i]: [int(j) for j in args.criteria[i]] for i in range(len(args.str_criteria))}


	# get predictions and ground truths
	with tables.open_file(args.pytable_path, 'r') as f:
		gts = np.array(f.root.labels[:]).flatten()
		preds = np.array(f.root.predictions[:]).flatten()
		print(gts.shape)
	
	# init presentation and compute coordinate grid.
	ppt = Presentation()
	grid_width = args.num_cols
	grid_height = args.num_rows
	grid = np.mgrid[0:grid_height, 0:grid_width]
	cartesian_coords = np.vstack([grid[0].ravel(), grid[1].ravel()]).T
	print(cartesian_coords)

	slidenames = criteria_dict.keys()
	for slidename in slidenames:
		criterion = criteria_dict[slidename]

		# find indices where the criterion is true
		inds = np.argwhere(np.logical_and(gts==criterion[0], preds==criterion[1])).flatten()
		grid_size = grid_width*grid_height
		num_imgs = min(grid_size, len(inds))
		selected_inds = np.random.choice(inds, num_imgs, replace=False)
		print(f'Number of images in {slidename}: {len(inds)}')

		# create new slide
		blank_slide_layout = ppt.slide_layouts[6]
		slide = ppt.slides.add_slide(blank_slide_layout)
		txBox = slide.shapes.add_textbox(Inches(ppt.slide_height.inches/2), Inches(0.2), Inches(2), Inches(1))	# add slide title
		txBox.text_frame.text = slidename

		# fill slide grid with selected images
		for j, ind in tqdm(enumerate(selected_inds)):
			coord = cartesian_coords[j]

			with tables.open_file(args.pytable_path, 'r') as f:	# we need to get the image_filename manually.
				img_filename = f.root.filenames[ind]
				label = f.root.labels[ind]
				img = f.root.imgs[ind, :, :, :]

			plt.imshow(img)
			title = img_filename.decode("utf-8").split('/')[-1]
			plt.title(title, wrap=True)
			plt.tick_params(color='b',bottom=False, left=False)
			plt.xticks([])
			plt.yticks([])
			with BytesIO() as img_buf:
				plt.savefig(img_buf, format='png', bbox_inches='tight')
				plt.close()
				im = Image.open(img_buf)
				addimagetoslide(slide, 
								img_buf, 
								top=Inches(coord[0]*1.2+1),
								left=Inches(coord[1]*1.2 + 0.25), 
								width=Inches(im.height/im.width), 
								height=Inches(1))
				im.close()

	ppt.save(args.ppt_save)
	
	